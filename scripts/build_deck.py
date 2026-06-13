"""Build a styled COOPUS pitch deck (.pptx) from a deck JSON file.

Usage:
    python scripts/build_deck.py <deck.json> <out.pptx>

The JSON shape (produced by the coopus-pitch-deck workflow):
{
  "deck_title": str,
  "deck_subtitle": str,
  "slides": [
    {"number": int, "title": str, "subtitle": str,
     "bullets": [str, ...], "stat_callout": str, "notes": str}, ...
  ]
}
"""
import json
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 16:9 canvas ---------------------------------------------------------
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

# ---- Palette (defence / tech) -------------------------------------------
NAVY      = RGBColor(0x0A, 0x0E, 0x1A)   # deep background
NAVY2     = RGBColor(0x10, 0x18, 0x2C)   # panel background
CYAN      = RGBColor(0x22, 0xD3, 0xEE)   # primary accent
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # hero-stat accent
WHITE     = RGBColor(0xF4, 0xF7, 0xFB)
MUTED     = RGBColor(0x9F, 0xB0, 0xC8)   # secondary text
RULE      = RGBColor(0x1E, 0x2A, 0x44)

FONT = "Segoe UI"


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=8, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
    return tb


def _bullet_box(slide, x, y, w, h, bullets, narrow=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    n = max(1, len(bullets))
    if n <= 3:
        size = 20
    elif n == 4:
        size = 19
    elif n == 5:
        size = 17
    else:
        size = 15
    if narrow:
        size -= 1
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        # accent tick
        tick = p.add_run()
        tick.text = "▸  "  # ▸
        tick.font.size = Pt(size)
        tick.font.color.rgb = CYAN
        tick.font.bold = True
        tick.font.name = FONT
        r = p.add_run()
        r.text = b
        r.font.size = Pt(size)
        r.font.color.rgb = WHITE
        r.font.name = FONT
    return tb


_UNIT = {"m/s", "km/h", "km", "m", "s", "%", "kg", "hz", "x", "×"}


def _split_stat(stat):
    """Return (hero, caption): a leading number (+unit) as the big figure, rest as caption."""
    tokens = stat.split(" ")
    hero = tokens[0]
    rest = tokens[1:]
    # pull a trailing unit token into the hero figure
    if rest and rest[0].lower() in _UNIT:
        hero = hero + " " + rest[0]
        rest = rest[1:]
    # if the very first token has no digit, don't force a split — show it all small-ish
    if not re.search(r"\d", tokens[0]):
        return stat, ""
    return hero, " ".join(rest)


def _notes(slide, text):
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def _footer(slide, idx, total, label="COOPUS — Cooperative Counter-UAS Defence"):
    _text(slide, Inches(0.55), Inches(7.02), Inches(9.0), Inches(0.4),
          [[(label, 10, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(11.6), Inches(7.02), Inches(1.2), Inches(0.4),
          [[(f"{idx:02d} / {total:02d}", 10, MUTED, False)]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    # accent band on the left
    _rect(s, 0, 0, Inches(0.25), EMU_H, CYAN)
    # small eyebrow
    _text(s, Inches(0.9), Inches(2.05), Inches(11), Inches(0.5),
          [[("COUNTER-UAS DEFENCE • VALIDATED IN SIMULATION", 14, CYAN, True)]])
    _text(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(2.2),
          [[(title, 54, WHITE, True)]], line_spacing=1.0)
    # rule
    _rect(s, Inches(0.95), Inches(4.55), Inches(3.2), Pt(3), AMBER)
    _text(s, Inches(0.9), Inches(4.8), Inches(11.2), Inches(1.6),
          [[(subtitle, 22, MUTED, False)]], line_spacing=1.1)
    _footer(s, 1, len(DECK["slides"]))
    return s


def content_slide(prs, idx, total, title, subtitle, bullets, stat):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    # header accent tick + title
    _rect(s, Inches(0.55), Inches(0.62), Inches(0.12), Inches(0.7), CYAN)
    _text(s, Inches(0.85), Inches(0.5), Inches(11.8), Inches(0.95),
          [[(title, 32, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    y = 1.55
    if subtitle:
        _text(s, Inches(0.85), Inches(1.5), Inches(11.8), Inches(0.5),
              [[(subtitle, 17, CYAN, False)]])
        y = 2.05
    # divider rule
    _rect(s, Inches(0.85), Inches(y), Inches(11.6), Pt(1.2), RULE)

    has_stat = bool(stat and stat.strip())
    bw = Inches(7.6) if has_stat else Inches(11.6)
    _bullet_box(s, Inches(0.85), Inches(y + 0.25), bw, Inches(4.4), bullets, narrow=has_stat)

    if has_stat:
        panel_x, panel_y, panel_w, panel_h = Inches(8.85), Inches(y + 0.25), Inches(3.6), Inches(4.0)
        _rect(s, panel_x, panel_y, panel_w, panel_h, NAVY2, line=RULE)
        big, rest = _split_stat(stat.strip())
        big_size = 60 if len(big) <= 4 else (46 if len(big) <= 7 else 36)
        _text(s, panel_x + Inches(0.15), panel_y + Inches(0.55), panel_w - Inches(0.3), Inches(1.7),
              [[(big, big_size, AMBER, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if rest:
            _text(s, panel_x + Inches(0.25), panel_y + Inches(2.35), panel_w - Inches(0.5), Inches(1.5),
                  [[(rest, 16, WHITE, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.1)
    _footer(s, idx, total)
    return s


def closing_slide(prs, idx, total, title, subtitle, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _rect(s, 0, 0, EMU_W, Inches(0.18), CYAN)
    _text(s, Inches(0.9), Inches(1.4), Inches(11.6), Inches(1.2),
          [[(title, 44, WHITE, True)]])
    _rect(s, Inches(0.95), Inches(2.55), Inches(3.2), Pt(3), AMBER)
    if subtitle:
        _text(s, Inches(0.9), Inches(2.8), Inches(11.2), Inches(0.9),
              [[(subtitle, 22, CYAN, False)]], line_spacing=1.1)
    if bullets:
        _bullet_box(s, Inches(0.9), Inches(3.9), Inches(11.4), Inches(2.6), bullets)
    _footer(s, idx, total)
    return s


def build(deck, out_path):
    global DECK
    DECK = deck
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    slides = deck["slides"]
    total = len(slides)
    for sl in slides:
        n = sl.get("number", 0)
        title = sl.get("title", "")
        subtitle = sl.get("subtitle", "") or ""
        bullets = sl.get("bullets", []) or []
        stat = sl.get("stat_callout", "") or ""
        if n == 1:
            s = title_slide(prs, deck.get("deck_title", title),
                            deck.get("deck_subtitle", subtitle) or subtitle)
        elif n == total:
            s = closing_slide(prs, n, total, title, subtitle, bullets)
        else:
            s = content_slide(prs, n, total, title, subtitle, bullets, stat)
        _notes(s, sl.get("notes", ""))
    prs.save(out_path)
    print(f"Wrote {out_path} ({total} slides)")


def _extract_deck(obj):
    """Accept a bare deck, a {deck: ...}, or a workflow {result: {deck: ...}}."""
    if "slides" in obj:
        return obj
    if "deck" in obj:
        return obj["deck"]
    if "result" in obj and isinstance(obj["result"], dict):
        return _extract_deck(obj["result"])
    raise ValueError("Could not find a deck (no 'slides' key) in the JSON")


if __name__ == "__main__":
    with open(sys.argv[1], "r", encoding="utf-8") as f:
        deck = _extract_deck(json.load(f))
    build(deck, sys.argv[2])
